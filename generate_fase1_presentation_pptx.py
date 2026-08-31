#!/usr/bin/env python3
"""
Generador Maestro de Presentación PPTX - Fase 1: Definición Proyecto APT
Asignatura: Portafolio de Título (APT122) - Duoc UC Sede San Joaquín
Proyecto: KRONOS SENTINEL - Autonomous AI-IPS & Real-Time Incident Voice Response SOAR Architecture
Tema: Lo-Fi Cyberpunk Purple Waves / Glassmorphic Deep Purple / 16:9 Widescreen
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# =============================================================================
# CONFIGURACIÓN GLOBAL & RUTAS
# =============================================================================
OUTPUT_DIR = "docs/Fase_1_Definicion_Proyecto_APT"
OUTPUT_PPTX = os.path.join(OUTPUT_DIR, "Presentacion_Proyecto_APT_Fase1_KRONOS_SENTINEL.pptx")
ASSETS_DIR = "assets"
BG_IMAGE_PATH = os.path.join(ASSETS_DIR, "lofi_purple_waves_bg.png")

# =============================================================================
# PALETA CROMÁTICA LO-FI PURPLE WAVES & SYNTH CYBER
# =============================================================================
C_CARD_BG       = RGBColor(20, 11, 41)       # #140B29 (Glassmorphic Deep Purple)
C_CARD_HEADER   = RGBColor(35, 20, 72)       # #231448 (Luminous Purple Header)
C_BORDER        = RGBColor(76, 29, 149)      # #4C1D95 (Neon Purple Border)
C_BORDER_LIGHT  = RGBColor(124, 58, 237)     # #7C3AED (Bright Violet Border)

# Acentos Neón / Lo-Fi Synth
C_CYAN          = RGBColor(0, 240, 255)      # #00F0FF (Electric Cyan)
C_SKY           = RGBColor(56, 189, 248)     # #38BDF8 (Sky Neon)
C_PURPLE        = RGBColor(192, 132, 252)    # #C084FC (Bright Lilac)
C_PINK          = RGBColor(244, 114, 182)    # #F472B6 (Synth Pink)
C_AMBER         = RGBColor(251, 191, 36)     # #FBBF24 (Amber Gold)
C_GREEN         = RGBColor(52, 211, 153)     # #34D399 (Emerald Green)
C_RED           = RGBColor(248, 113, 113)    # #F87171 (Signal Coral Red)

# Escala Tipográfica
C_TEXT_WHITE    = RGBColor(255, 255, 255)    # #FFFFFF
C_TEXT_LIGHT    = RGBColor(237, 233, 254)    # #EDE9FE (Soft Lavender White)
C_TEXT_MUTED    = RGBColor(196, 181, 253)    # #C4B5FD (Muted Lilac)
C_TEXT_DARK     = RGBColor(139, 92, 246)     # #8B5CF6 (Deep Violet Footer)

# =============================================================================
# HELPER FUNCTIONS - COMPONENTES VISUALES LO-FI
# =============================================================================
def apply_lofi_background(slide):
    """Aplica el fondo Lo-Fi Purple Waves 16:9 widescreen."""
    if os.path.exists(BG_IMAGE_PATH):
        slide.shapes.add_picture(BG_IMAGE_PATH, 0, 0, Inches(13.333), Inches(7.5))
    else:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(15, 7, 30)
        bg.line.fill.background()

def add_header(slide, title_text, category_text="PORTAFOLIO DE TÍTULO (APT122) • FASE 1: DEFINICIÓN"):
    """Renderiza el banner superior temático con acento Neón Cyan/Purple."""
    tb_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.733), Inches(0.25))
    tf_cat = tb_cat.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = "Consolas"
    p_cat.font.size = Pt(8.5)
    p_cat.font.bold = True
    p_cat.font.color.rgb = C_CYAN

    tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.58), Inches(11.733), Inches(0.50))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(19)
    p_title.font.bold = True
    p_title.font.color.rgb = C_TEXT_WHITE

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(11.733), Inches(0.015)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C_BORDER_LIGHT
    line.line.fill.background()

def add_footer(slide, current_page, total_pages=12):
    """Pie de página institucional con branding Lo-Fi."""
    tb_ft = slide.shapes.add_textbox(Inches(0.8), Inches(7.10), Inches(11.733), Inches(0.25))
    tf_ft = tb_ft.text_frame
    tf_ft.word_wrap = True
    p_ft = tf_ft.paragraphs[0]
    p_ft.text = f"DUOC UC SAN JOAQUÍN • ESCUELA DE INFORMÁTICA Y TELECOMUNICACIONES  |  KRONOS SENTINEL  |  SLIDE {current_page:02d} / {total_pages:02d}"
    p_ft.font.name = "Consolas"
    p_ft.font.size = Pt(8)
    p_ft.font.color.rgb = C_TEXT_DARK

def add_card(slide, x, y, w, h, title="", border_color=C_BORDER, fill_color=C_CARD_BG, header_color=C_CARD_HEADER, header_h=0.38):
    """Crea una tarjeta glassmorphic púrpura con esquinas redondeadas."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.0)

    if title:
        hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(header_h))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = header_color
        hdr.line.color.rgb = border_color
        hdr.line.width = Pt(1.0)

        tf = hdr.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        tf.margin_top = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = C_CYAN

    return card

def add_stat_card(slide, x, y, w, h, metric, label, desc="", color=C_CYAN):
    """Crea una micro-tarjeta de estadística con resplandor neón."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = C_CARD_BG
    card.line.color.rgb = color
    card.line.width = Pt(1.2)

    tb = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.08), Inches(w - 0.2), Inches(h - 0.16))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    p1 = tf.paragraphs[0]
    p1.text = metric
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = color

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(9.5)
    p2.font.bold = True
    p2.font.color.rgb = C_TEXT_WHITE

    if desc:
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.name = "Segoe UI"
        p3.font.size = Pt(8)
        p3.font.color.rgb = C_TEXT_MUTED

def add_pill_badge(slide, x, y, w, h, text, text_color=C_CYAN, bg_color=C_CARD_HEADER, border_color=C_BORDER_LIGHT, font_size=8.5):
    """Crea una píldora compacta para tags y estados."""
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    pill.fill.solid()
    pill.fill.fore_color.rgb = bg_color
    pill.line.color.rgb = border_color
    pill.line.width = Pt(0.8)

    tf = pill.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Consolas"
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = text_color
    return pill

def add_image_pro(slide, image_path, x, y, w, h):
    """Inserta una imagen manteniendo su encuadre armónico."""
    if os.path.exists(image_path):
        return slide.shapes.add_picture(image_path, Inches(x), Inches(y), Inches(w), Inches(h))
    return None

# =============================================================================
# CONSTRUCTOR PRINCIPAL DE LAS 12 DIAPOSITIVAS
# =============================================================================
def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # -------------------------------------------------------------------------
    # SLIDE 1: PORTADA EJECUTIVA (Lo-Fi Ambient)
    # -------------------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    apply_lofi_background(s1)

    add_pill_badge(s1, 0.8, 0.65, 4.4, 0.30, "PORTAFOLIO DE TÍTULO (APT122) • FASE 1", text_color=C_CYAN, border_color=C_CYAN)
    add_pill_badge(s1, 5.3, 0.65, 2.5, 0.30, "DEFENSA EN PROFUNDIDAD", text_color=C_GREEN, border_color=C_GREEN)

    tb_hero = s1.shapes.add_textbox(Inches(0.8), Inches(1.10), Inches(8.5), Inches(1.9))
    tf_hero = tb_hero.text_frame
    tf_hero.word_wrap = True

    p1 = tf_hero.paragraphs[0]
    p1.text = "KRONOS SENTINEL"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = C_TEXT_WHITE

    p2 = tf_hero.add_paragraph()
    p2.text = "Autonomous AI-IPS & Real-Time Incident Voice Response SOAR Architecture"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(15)
    p2.font.bold = True
    p2.font.color.rgb = C_SKY

    p3 = tf_hero.add_paragraph()
    p3.text = "Mitigación Atómica en Kernel FreeBSD + Orquestación de Voz en Tiempo Real con IA Multimodal"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(10.5)
    p3.font.color.rgb = C_TEXT_MUTED

    # Emblema Escudo (Con fondo dark púrpura integrado)
    logo_path = os.path.join(ASSETS_DIR, "sentinel_shield_logo.png")
    add_image_pro(s1, logo_path, 9.8, 0.75, 2.4, 2.4)

    # Tarjeta Izquierda: Datos Académicos
    add_card(s1, 0.8, 3.25, 5.7, 3.55, title="DATOS ACADÉMICOS E INSTITUCIONALES", border_color=C_BORDER_LIGHT)
    tb_acad = s1.shapes.add_textbox(Inches(1.0), Inches(3.75), Inches(5.3), Inches(2.9))
    tf_acad = tb_acad.text_frame
    tf_acad.word_wrap = True

    acad_items = [
        ("Institución:", "Duoc UC — Sede San Joaquín"),
        ("Escuela:", "Escuela de Informática y Telecomunicaciones"),
        ("Carrera:", "Ingeniería en Conectividad y Redes"),
        ("Asignatura:", "Portafolio de Título (APT122)"),
        ("Estudiante:", "Bruno Urrea Ortiz (RUT: 21.543.637-3)"),
        ("Fecha Exposición:", "8 de Septiembre de 2026")
    ]
    for idx, (lbl, val) in enumerate(acad_items):
        p = tf_acad.paragraphs[0] if idx == 0 else tf_acad.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{lbl} "
        r1.font.name = "Segoe UI"
        r1.font.bold = True
        r1.font.size = Pt(9.5)
        r1.font.color.rgb = C_CYAN

        r2 = p.add_run()
        r2.text = val
        r2.font.name = "Segoe UI"
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = C_TEXT_LIGHT

    # Tarjeta Derecha: Equipo de Ingeniería
    add_card(s1, 6.833, 3.25, 5.7, 3.55, title="EQUIPO DE INGENIERÍA Y ROLES ASIGNADOS", border_color=C_BORDER_LIGHT)
    tb_team = s1.shapes.add_textbox(Inches(7.033), Inches(3.75), Inches(5.3), Inches(2.9))
    tf_team = tb_team.text_frame
    tf_team.word_wrap = True

    team_members = [
        ("Bruno Urrea Ortiz", "Líder de Ciberseguridad, Motor KRONOS y Gemini Live", C_CYAN),
        ("Freddy Vásquez Cortés", "Ingeniería de Routing, Switching L2/L3 y Telefonía Asterisk", C_SKY),
        ("Cristóbal Quezada", "Administración de Servicios Web, Proxy HAProxy y DMZ", C_AMBER),
        ("Kevin Retamales", "Hardening Perimetral, Inteligencia pfBlockerNG y QA", C_GREEN)
    ]
    for idx, (name, role, col) in enumerate(team_members):
        p = tf_team.paragraphs[0] if idx == 0 else tf_team.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {name}\n"
        r1.font.name = "Segoe UI"
        r1.font.bold = True
        r1.font.size = Pt(9.5)
        r1.font.color.rgb = C_TEXT_WHITE

        r2 = p.add_run()
        r2.text = f"   {role}\n"
        r2.font.name = "Segoe UI"
        r2.font.size = Pt(8.5)
        r2.font.color.rgb = col

    add_footer(s1, 1)

    # -------------------------------------------------------------------------
    # SLIDE 2: AGENDA DE LA PRESENTACIÓN
    # -------------------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    apply_lofi_background(s2)
    add_header(s2, "Agenda de la Presentación (10 Secciones)")

    sections_left = [
        ("01", "Contexto y Problemática", "Crisis de falsos positivos y colapso cognitivo del operador.", C_RED),
        ("02", "Descripción del Proyecto", "Defensa en profundidad y arquitectura SOAR a costo $0 CLP.", C_CYAN),
        ("03", "Vínculo con el Perfil de Egreso", "Mapeo de 6 competencias troncales de Conectividad y Redes.", C_SKY),
        ("04", "Intereses Profesionales", "Especialización en Ciberseguridad, SecOps y trayecto DevSec.", C_PURPLE),
        ("05", "Factibilidad y Gestión de Riesgos", "Semestre académico, recursos Open Source y evasión de CGNAT.", C_AMBER)
    ]
    sections_right = [
        ("06", "Objetivos General y Específicos", "8 metas técnicas medibles de ingeniería y arquitectura.", C_CYAN),
        ("07", "Metodología y Roles del Equipo", "Ciclo iterativo e incremental de 4 fases operativas.", C_GREEN),
        ("08", "Evidencias Comprometidas", "Entregables de avance, código IaC y demostración en vivo.", C_SKY),
        ("09", "Plan de Trabajo y Carta Gantt", "Cronograma de 18 semanas académicas y dependencias.", C_AMBER),
        ("10", "Cierre y Próximos Pasos (Fase 2)", "Hitos de implementación perimetral y ronda de preguntas.", C_GREEN)
    ]

    for idx, (num, title, desc, col_acc) in enumerate(sections_left):
        y = 1.35 + idx * 1.10
        add_card(s2, 0.8, y, 5.7, 0.95, border_color=C_BORDER)
        add_pill_badge(s2, 0.95, y + 0.15, 0.55, 0.28, num, text_color=col_acc, bg_color=C_CARD_HEADER, border_color=col_acc)

        tb = s2.shapes.add_textbox(Inches(1.6), Inches(y + 0.10), Inches(4.75), Inches(0.75))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{title}\n"
        r1.font.name = "Segoe UI"
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = C_TEXT_WHITE

        r2 = p.add_run()
        r2.text = desc
        r2.font.name = "Segoe UI"
        r2.font.size = Pt(8.5)
        r2.font.color.rgb = C_TEXT_MUTED

    for idx, (num, title, desc, col_acc) in enumerate(sections_right):
        y = 1.35 + idx * 1.10
        add_card(s2, 6.833, y, 5.7, 0.95, border_color=C_BORDER)
        add_pill_badge(s2, 6.983, y + 0.15, 0.55, 0.28, num, text_color=col_acc, bg_color=C_CARD_HEADER, border_color=col_acc)

        tb = s2.shapes.add_textbox(Inches(7.633), Inches(y + 0.10), Inches(4.75), Inches(0.75))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{title}\n"
        r1.font.name = "Segoe UI"
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = C_TEXT_WHITE

        r2 = p.add_run()
        r2.text = desc
        r2.font.name = "Segoe UI"
        r2.font.size = Pt(8.5)
        r2.font.color.rgb = C_TEXT_MUTED

    add_footer(s2, 2)

    # -------------------------------------------------------------------------
    # SLIDE 3: CONTEXTO Y PROBLEMÁTICA (Métricas & Diagrama)
    # -------------------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    apply_lofi_background(s3)
    add_header(s3, "Contexto y Problemática en la Industria de Ciberseguridad")

    # Métricas Visuales Superiores
    add_stat_card(s3, 0.8, 1.35, 2.7, 1.55, ">50%", "Falsos Positivos", "Fatiga extrema en analistas SOC N1/N2 por firmas genéricas.", C_RED)
    add_stat_card(s3, 3.7, 1.35, 2.7, 1.55, "0.05ms", "Visión de Túnel", "El humano colapsa al intentar aislar firewall y reportar al CISO.", C_AMBER)

    # Tarjeta de Solución Autónoma
    add_card(s3, 0.8, 3.05, 5.6, 2.65, title="PROPUESTA DE VALOR: KRONOS SENTINEL", border_color=C_GREEN)
    tb_sol = s3.shapes.add_textbox(Inches(0.95), Inches(3.50), Inches(5.3), Inches(2.1))
    tf_sol = tb_sol.text_frame
    tf_sol.word_wrap = True
    p_sol = tf_sol.paragraphs[0]
    p_sol.text = (
        "• Extracción del Factor Humano: Desacopla la contención atómica en kernel de la notificación.\n"
        "• Doble Acción Simultánea:\n"
        "   1. Drop instantáneo en FreeBSD vía pfctl -k y tabla <snort2c> (<100ms).\n"
        "   2. Llamada interactiva por voz mediante Asterisk y Gemini Live (<1.5s)."
    )
    p_sol.font.name = "Segoe UI"
    p_sol.font.size = Pt(8.8)
    p_sol.font.color.rgb = C_TEXT_LIGHT

    # Diagrama de Mitigación en Kernel
    flow_path = os.path.join(ASSETS_DIR, "pfctl_decision_flow.png")
    add_image_pro(s3, flow_path, 6.70, 1.35, 5.833, 3.30)

    # Banner Inferior: Ámbito de Impacto
    add_card(s3, 0.8, 5.85, 11.733, 1.05, title="", border_color=C_BORDER_LIGHT, fill_color=C_CARD_HEADER)
    tb_amb = s3.shapes.add_textbox(Inches(1.0), Inches(5.92), Inches(11.333), Inches(0.9))
    tf_amb = tb_amb.text_frame
    tf_amb.word_wrap = True
    p_amb = tf_amb.paragraphs[0]
    r1 = p_amb.add_run()
    r1.text = "ÁMBITO & POBLACIÓN DE IMPACTO: "
    r1.font.name = "Consolas"
    r1.font.bold = True
    r1.font.size = Pt(9)
    r1.font.color.rgb = C_CYAN

    r2 = p_amb.add_run()
    r2.text = "Datacenters, plataformas web expuestas e infraestructuras corporativas en Chile y LATAM. Beneficia directamente a Analistas SOC, CISOs y Administradores de Redes."
    r2.font.name = "Segoe UI"
    r2.font.size = Pt(8.5)
    r2.font.color.rgb = C_TEXT_LIGHT

    add_footer(s3, 3)

    # -------------------------------------------------------------------------
    # SLIDE 4: DESCRIPCIÓN DEL PROYECTO (Topología Visual & 4 Pilares)
    # -------------------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    apply_lofi_background(s4)
    add_header(s4, "Descripción del Proyecto KRONOS SENTINEL")

    # Banner Misión
    add_card(s4, 0.8, 1.35, 11.733, 0.75, title="", border_color=C_BORDER_LIGHT, fill_color=C_CARD_HEADER)
    tb_m = s4.shapes.add_textbox(Inches(1.0), Inches(1.40), Inches(11.333), Inches(0.65))
    tf_m = tb_m.text_frame
    tf_m.word_wrap = True
    p_m = tf_m.paragraphs[0]
    r1 = p_m.add_run()
    r1.text = "OBJETIVO CENTRAL ($0 CLP): "
    r1.font.name = "Consolas"
    r1.font.bold = True
    r1.font.size = Pt(9)
    r1.font.color.rgb = C_CYAN

    r2 = p_m.add_run()
    r2.text = "“Arquitectura SOAR autónoma de costo cero que contiene intrusiones en kernel FreeBSD en microsegundos e interactúa por voz en tiempo real con el CISO mediante Inteligencia Artificial.”"
    r2.font.name = "Segoe UI"
    r2.font.size = Pt(9.2)
    r2.font.bold = True
    r2.font.color.rgb = C_TEXT_WHITE

    # Topología Global
    arch_path = os.path.join(ASSETS_DIR, "architecture_diagram.png")
    add_image_pro(s4, arch_path, 0.8, 2.25, 5.80, 3.15)

    # Widget Inferior Bajo la Topología
    add_card(s4, 0.8, 5.55, 5.80, 1.35, title="DEFENSA MULTI-CAPA ZERO TRUST", border_color=C_GREEN, header_h=0.32)
    tb_zt = s4.shapes.add_textbox(Inches(0.95), Inches(5.90), Inches(5.50), Inches(0.95))
    tf_zt = tb_zt.text_frame
    tf_zt.word_wrap = True
    p_zt = tf_zt.paragraphs[0]
    p_zt.text = "• 4 VLANs 802.1Q aisladas (Corp 10, DMZ 20, VoIP 30, Mgmt 99).\n• Aislamiento estricto de DMZ con proxy inverso HAProxy y terminación HTTPS."
    p_zt.font.name = "Segoe UI"
    p_zt.font.size = Pt(8.2)
    p_zt.font.color.rgb = C_TEXT_LIGHT

    # 4 Pilares Modulares
    pilars_s4 = [
        ("01. PERÍMETRO INLINE IPS", "pfSense CE 2.9.0 + Suricata 7.x Inline Netmap (Hardware Drop ring-buffer) y pfBlockerNG GeoIP MaxMind.", C_CYAN),
        ("02. CAPA WEB DMZ & PROXY", "HAProxy 2.8+ SSL y Stick-Tables anti-fuzzing L7 protegiendo servidor vulnerable DVWA en VLAN 20.", C_GREEN),
        ("03. MOTOR KRONOS AST", "Python 3.12 AST parser (supresión >50% ruido) y control de kernel FreeBSD pfctl (Kill States y <snort2c>).", C_AMBER),
        ("04. SOAR VOZ GEMINI & PBX", "Asterisk 20 LTS en Docker con auto-dialer AMI y Google Gemini Live Flash 3.1 para debriefing hablado.", C_PURPLE)
    ]
    for idx, (title, desc, col) in enumerate(pilars_s4):
        y_p = 2.25 + idx * 1.16
        add_card(s4, 6.833, y_p, 5.70, 1.08, title=title, border_color=col, header_h=0.32)
        tb_p = s4.shapes.add_textbox(Inches(6.983), Inches(y_p + 0.36), Inches(5.40), Inches(0.68))
        tf_p = tb_p.text_frame
        tf_p.word_wrap = True
        p_p = tf_p.paragraphs[0]
        p_p.text = desc
        p_p.font.name = "Segoe UI"
        p_p.font.size = Pt(8.2)
        p_p.font.color.rgb = C_TEXT_LIGHT

    add_footer(s4, 4)

    # -------------------------------------------------------------------------
    # SLIDE 5: VÍNCULO CON EL PERFIL DE EGRESO (6 Competencias)
    # -------------------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    apply_lofi_background(s5)
    add_header(s5, "Vínculo con el Perfil de Egreso (6 Competencias Troncales APT)")

    comps = [
        ("Comp. 8: Planes de Prevención y Respuesta a Riesgos", "Ciberseguridad Defensiva", "Diseño de topología perimetral en pfSense CE 2.9.0, reglas Zero Trust, Suricata Inline Netmap IPS y mitigación atómica de ataques en kernel con pfctl.", C_CYAN),
        ("Comp. 6: Automatizar Procesos y Gestión de Red", "Automatización & SecOps", "Desarrollo en Python 3.12 del Motor de Correlación KRONOS, parsing estructurado EVE JSON, análisis sintáctico AST y despacho asíncrono de webhooks.", C_SKY),
        ("Comp. 5: Unificar Servicios de Voz, Datos y Video", "Telecomunicaciones & VoIP", "Despliegue de centralita Asterisk 20 LTS en Docker, troncales PJSIP, configuración de Dialplans y conexión RTP streaming con agentes de voz.", C_GREEN),
        ("Comp. 4: Controlar y Operar Redes Corporativas", "Infraestructura de Redes", "Segmentación L2/L3 en 4 VLANs 802.1Q (Corp 10, DMZ 20, VoIP 30, Mgmt 99), direccionamiento CIDR, DHCP relay y conmutación virtualizada.", C_AMBER),
        ("Comp. 3: Adaptar Tecnologías de Punta y Tendencias", "Innovación & Ciberdefensa", "Integración de IA Generativa Multimodal en tiempo real (Google Gemini Live Flash 3.1) y túneles Zero Trust Mesh con Tailscale Subnet Router.", C_PURPLE),
        ("Comp. 7: Gestionar Seguridad ante Vulnerabilidades", "Seguridad Ofensiva/Defensiva", "Análisis de vectores de ataque web (SQLi, RCE, Fuzzing), despliegue de laboratorio DVWA y hardening de HAProxy con Stick-Tables dinámicas en RAM.", C_RED)
    ]
    card_w5 = 5.74
    card_h5 = 1.65
    for idx, (title, area, desc, color) in enumerate(comps):
        col = idx % 2
        row = idx // 2
        x = 0.8 + col * (card_w5 + 0.25)
        y = 1.35 + row * (card_h5 + 0.18)

        add_card(s5, x, y, card_w5, card_h5, border_color=color)
        tb = s5.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.10), Inches(card_w5 - 0.3), Inches(card_h5 - 0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = f"{title}\n"
        r1.font.name = "Segoe UI"
        r1.font.bold = True
        r1.font.size = Pt(9.2)
        r1.font.color.rgb = color

        r2 = p1.add_run()
        r2.text = f"Área de Desempeño: {area}\n"
        r2.font.name = "Segoe UI"
        r2.font.bold = True
        r2.font.size = Pt(8.2)
        r2.font.color.rgb = C_TEXT_MUTED

        r3 = p1.add_run()
        r3.text = desc
        r3.font.name = "Segoe UI"
        r3.font.size = Pt(8.0)
        r3.font.color.rgb = C_TEXT_LIGHT

    add_footer(s5, 5)

    # -------------------------------------------------------------------------
    # SLIDE 6: INTERESES Y PROYECCIÓN PROFESIONAL
    # -------------------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    apply_lofi_background(s6)
    add_header(s6, "Relación con Intereses y Proyección Profesional")

    # Tarjeta Izquierda
    add_card(s6, 0.8, 1.35, 5.7, 5.55, title="ESPECIALIZACIÓN EN CIBERSEGURIDAD Y SECOPS", border_color=C_CYAN)
    tb_m = s6.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(5.3), Inches(4.9))
    tf_m = tb_m.text_frame
    tf_m.word_wrap = True
    m_text = (
        "• Pasión por la Arquitectura Defensiva:\n"
        "Consolidación de competencias en Hardening Perimetral, Prevención de Intrusiones (IPS) y Sistemas Autónomos SOAR.\n\n"
        "• Experiencia Práctica y Trayectoria:\n"
        "Aprendizajes aplicados en Home Labs Proxmox VE, administración de firewalls pfSense en producción y liderazgo en torneos CTF con el equipo DevSec.\n\n"
        "• Proyección Laboral a 5 Años:\n"
        "Desempeñarse como Ingeniero de Ciberseguridad / Arquitecto SecOps liderando equipos de respuesta ante incidentes (CSIRT/SOC)."
    )
    p = tf_m.paragraphs[0]
    p.text = m_text
    p.font.name = "Segoe UI"
    p.font.size = Pt(9.2)
    p.font.color.rgb = C_TEXT_LIGHT

    # Tarjeta Derecha
    add_card(s6, 6.833, 1.35, 5.7, 5.55, title="APORTE AL DESARROLLO DEL EQUIPO DUOC UC", border_color=C_SKY)
    tb_d = s6.shapes.add_textbox(Inches(7.033), Inches(1.85), Inches(5.3), Inches(4.9))
    tf_d = tb_d.text_frame
    tf_d.word_wrap = True
    d_text = (
        "• Dominio de Entornos Enterprise $0 CLP:\n"
        "Capacidad de implementar arquitecturas de seguridad de nivel corporativo utilizando 100% tecnologías Open Source sin costos de licenciamiento.\n\n"
        "• Interoperabilidad Multi-Disciplina:\n"
        "Convergencia real entre Redes L2/L3 (VLANs 802.1Q), Telefonía IP empresarial (Asterisk PBX), Hardening Unix (FreeBSD) e Inteligencia Artificial.\n\n"
        "• Alineación con Estándares Industriales:\n"
        "Alineación directa con certificaciones Cisco CCNA (200-301), CompTIA Security+ y marco Zero Trust NIST SP 800-207."
    )
    p = tf_d.paragraphs[0]
    p.text = d_text
    p.font.name = "Segoe UI"
    p.font.size = Pt(9.2)
    p.font.color.rgb = C_TEXT_LIGHT

    add_footer(s6, 6)

    # -------------------------------------------------------------------------
    # SLIDE 7: FACTIBILIDAD Y GESTIÓN DE RIESGOS
    # -------------------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    apply_lofi_background(s7)
    add_header(s7, "Factibilidad Técnica, Recursos y Gestión de Riesgos")

    # Métricas Superiores de Factibilidad
    add_stat_card(s7, 0.8, 1.35, 2.7, 1.45, "$0 CLP", "100% Open Source", "FreeBSD 14, pfSense, Suricata, Docker, Asterisk.", C_GREEN)
    add_stat_card(s7, 3.7, 1.35, 2.7, 1.45, "18 Sem.", "Temporalidad", "72h taller presencial + 144h laboratorio autónomo.", C_CYAN)

    # Tarjeta de Infraestructura
    add_card(s7, 0.8, 2.95, 5.6, 3.95, title="INFRAESTRUCTURA & CAPAS GRATUITAS", border_color=C_GREEN)
    tb_inf = s7.shapes.add_textbox(Inches(0.95), Inches(3.40), Inches(5.3), Inches(3.3))
    tf_inf = tb_inf.text_frame
    tf_inf.word_wrap = True
    p_inf = tf_inf.paragraphs[0]
    p_inf.text = (
        "• Infraestructura de Redes:\n"
        "   - Laboratorios Duoc UC San Joaquín y plataforma Proxmox VE.\n\n"
        "• Capas Gratuitas Comunitarias:\n"
        "   - Google AI Studio (Gemini Live Flash 3.1 Free Tier).\n"
        "   - MaxMind GeoLite2 (Listas GeoIP actualizadas).\n\n"
        "• Validación Presupuestaria:\n"
        "   - Cero inversión en hardware privativo o licencias comerciales."
    )
    p_inf.font.name = "Segoe UI"
    p_inf.font.size = Pt(8.8)
    p_inf.font.color.rgb = C_TEXT_LIGHT

    # Tarjeta Derecha: Matriz de Riesgos & Mitigaciones
    add_card(s7, 6.833, 1.35, 5.7, 5.55, title="MATRIZ DE RIESGOS Y MITIGACIONES TÉCNICAS", border_color=C_AMBER)
    tb_r = s7.shapes.add_textbox(Inches(7.033), Inches(1.85), Inches(5.3), Inches(4.9))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    r_text = (
        "🔴 Riesgo 1: Bloqueo de Puertos / CGNAT en Redes Remotas\n"
        "   -> Mitigación: Despliegue de Tailscale Subnet Router. Publica la VLAN VoIP vía túnel WireGuard Zero Trust sin requerir IP pública.\n\n"
        "🔴 Riesgo 2: Falsos Positivos Saturando Llamadas al CISO\n"
        "   -> Mitigación: Algoritmo heurístico AST en Python 3.12 que valida sintaxis SQLi real y exige confirmación en tabla <snort2c>.\n\n"
        "🟠 Riesgo 3: Latencia en la Interacción de Voz de IA\n"
        "   -> Mitigación: Conexión directa por WebSocket bidireccional (PCM 24kHz) hacia Gemini Live Flash 3.1 (<400ms)."
    )
    p = tf_r.paragraphs[0]
    p.text = r_text
    p.font.name = "Segoe UI"
    p.font.size = Pt(8.8)
    p.font.color.rgb = C_TEXT_LIGHT

    add_footer(s7, 7)

    # -------------------------------------------------------------------------
    # SLIDE 8: OBJETIVOS GENERAL Y ESPECÍFICOS
    # -------------------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    apply_lofi_background(s8)
    add_header(s8, "Objetivos General y Específicos del Proyecto")

    # Objetivo General
    add_card(s8, 0.8, 1.35, 11.733, 1.25, title="OBJETIVO GENERAL DEL PROYECTO", border_color=C_CYAN)
    add_pill_badge(s8, 10.1, 1.40, 2.2, 0.28, "META ARQUITECTÓNICA", text_color=C_CYAN, border_color=C_CYAN, font_size=8)

    tb_og = s8.shapes.add_textbox(Inches(1.0), Inches(1.75), Inches(11.333), Inches(0.75))
    tf_og = tb_og.text_frame
    tf_og.word_wrap = True
    p_og = tf_og.paragraphs[0]
    p_og.text = "Diseñar, implementar y validar una arquitectura de defensa en profundidad y respuesta autónoma ante incidentes (SOAR) de costo cero ($0 CLP) denominada KRONOS SENTINEL, integrando prevención de intrusiones en kernel, supresión heurística de falsos positivos e interacción de voz en tiempo real con Inteligencia Artificial hacia el CISO."
    p_og.font.name = "Segoe UI"
    p_og.font.size = Pt(9.5)
    p_og.font.bold = True
    p_og.font.color.rgb = C_TEXT_WHITE

    # Objetivos Específicos (1 - 4)
    add_card(s8, 0.8, 2.75, 5.7, 4.15, title="OBJETIVOS ESPECÍFICOS (1 A 4)", border_color=C_SKY)
    tb_o1 = s8.shapes.add_textbox(Inches(1.0), Inches(3.20), Inches(5.3), Inches(3.5))
    tf_o1 = tb_o1.text_frame
    tf_o1.word_wrap = True
    objs_left = [
        ("OBJ-01", "Segmentación L2/L3", "Implementar 4 VLANs 802.1Q (Corp 10, DMZ 20, VoIP 30, Mgmt 99) en pfSense CE 2.9.0."),
        ("OBJ-02", "Prevención Netmap IPS", "Configurar Suricata 7.x en modo Inline IPS con descarte en hardware ring-buffer."),
        ("OBJ-03", "Capa DMZ Segura", "Desplegar HAProxy 2.8+ SSL y Stick-Tables anti-fuzzing L7 protegiendo servidor DVWA."),
        ("OBJ-04", "Motor KRONOS AST", "Desarrollar en Python 3.12 el analizador sintáctico para suprimir >50% de falsos positivos.")
    ]
    for idx, (tag, name, desc) in enumerate(objs_left):
        p = tf_o1.paragraphs[0] if idx == 0 else tf_o1.add_paragraph()
        r1 = p.add_run()
        r1.text = f"[{tag}] {name}: "
        r1.font.name = "Consolas"
        r1.font.bold = True
        r1.font.size = Pt(9)
        r1.font.color.rgb = C_SKY

        r2 = p.add_run()
        r2.text = f"{desc}\n\n"
        r2.font.name = "Segoe UI"
        r2.font.size = Pt(8.3)
        r2.font.color.rgb = C_TEXT_LIGHT

    # Objetivos Específicos (5 - 8)
    add_card(s8, 6.833, 2.75, 5.7, 4.15, title="OBJETIVOS ESPECÍFICOS (5 A 8)", border_color=C_GREEN)
    tb_o2 = s8.shapes.add_textbox(Inches(7.033), Inches(3.20), Inches(5.3), Inches(3.5))
    tf_o2 = tb_o2.text_frame
    tf_o2.word_wrap = True
    objs_right = [
        ("OBJ-05", "Kernel FreeBSD pfctl", "Integrar terminación de estados (pfctl -k) y verificación atómica en tabla <snort2c>."),
        ("OBJ-06", "Telefonía Asterisk PBX", "Desplegar Asterisk 20 LTS en Docker con canal PJSIP y auto-dialer AMI al softphone."),
        ("OBJ-07", "Agente de Voz Gemini Live", "Construir cliente WebSocket (PCM 24kHz) para debriefing táctico hablado en tiempo real."),
        ("OBJ-08", "Validación QA Integral", "Ejecutar matrices de penetración y latencia total de respuesta <1.5 segundos.")
    ]
    for idx, (tag, name, desc) in enumerate(objs_right):
        p = tf_o2.paragraphs[0] if idx == 0 else tf_o2.add_paragraph()
        r1 = p.add_run()
        r1.text = f"[{tag}] {name}: "
        r1.font.name = "Consolas"
        r1.font.bold = True
        r1.font.size = Pt(9)
        r1.font.color.rgb = C_GREEN

        r2 = p.add_run()
        r2.text = f"{desc}\n\n"
        r2.font.name = "Segoe UI"
        r2.font.size = Pt(8.3)
        r2.font.color.rgb = C_TEXT_LIGHT

    add_footer(s8, 8)

    # -------------------------------------------------------------------------
    # SLIDE 9: METODOLOGÍA Y ROLES DEL EQUIPO
    # -------------------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    apply_lofi_background(s9)
    add_header(s9, "Metodología de Trabajo y Roles del Equipo")

    etapas = [
        ("FASE 1: DISEÑO & TOPOLOGÍA", "Modelamiento IP, VLANs 802.1Q y tuning de pfSense CE 2.9.0.", C_CYAN),
        ("FASE 2: PERÍMETRO & DMZ", "Suricata Netmap IPS, pfBlockerNG GeoIP y HAProxy SSL DMZ.", C_SKY),
        ("FASE 3: KRONOS & VOZ IA", "Motor Python AST, wrappers pfctl FreeBSD, Asterisk y Gemini Live.", C_AMBER),
        ("FASE 4: INTEGRACIÓN & QA", "Simulación SQLi en vivo, latencia (<1.5s) y documentación formal.", C_GREEN)
    ]
    card_w9 = 2.78
    for idx, (title, desc, color) in enumerate(etapas):
        x = 0.8 + idx * (card_w9 + 0.20)
        add_card(s9, x, 1.35, card_w9, 1.70, title=title, border_color=color, header_h=0.35)
        tb = s9.shapes.add_textbox(Inches(x + 0.12), Inches(1.78), Inches(card_w9 - 0.24), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = "Segoe UI"
        p.font.size = Pt(8.5)
        p.font.color.rgb = C_TEXT_LIGHT

    roles = [
        ("Bruno Urrea Ortiz", "Líder de Ciberseguridad & Motor KRONOS", "Arquitectura global, log_correlator.py, algoritmo AST anti-ruido, control en kernel pfctl y cliente Gemini Live.", C_CYAN),
        ("Freddy Vásquez Cortés", "Ingeniería de Routing, Switching & VoIP", "Troncal 802.1Q en pfSense, centralita Asterisk 20 LTS Docker, auto-dialer AMI y Tailscale Mesh.", C_SKY),
        ("Cristóbal Quezada", "Servicios Web & Proxy Inverso DMZ", "HAProxy 2.8+ SSL, Stick-Tables anti-fuzzing y administración de DVWA en VLAN 20.", C_AMBER),
        ("Kevin Retamales", "Hardening Perimetral, GeoIP & QA Lead", "Políticas Zero Trust en pfSense, pfBlockerNG-devel GeoIP y diseño de matrices QA.", C_GREEN)
    ]
    for idx, (name, role, desc, color) in enumerate(roles):
        x = 0.8 + idx * (card_w9 + 0.20)
        add_card(s9, x, 3.25, card_w9, 3.65, title=f"ROL: {name.split()[0].upper()}", border_color=color, header_h=0.35)
        tb = s9.shapes.add_textbox(Inches(x + 0.12), Inches(3.70), Inches(card_w9 - 0.24), Inches(3.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = f"{name}\n"
        r1.font.name = "Segoe UI"
        r1.font.bold = True
        r1.font.size = Pt(9.5)
        r1.font.color.rgb = C_TEXT_WHITE

        r2 = p1.add_run()
        r2.text = f"{role}\n\n"
        r2.font.name = "Segoe UI"
        r2.font.bold = True
        r2.font.size = Pt(8.5)
        r2.font.color.rgb = color

        r3 = p1.add_run()
        r3.text = desc
        r3.font.name = "Segoe UI"
        r3.font.size = Pt(8)
        r3.font.color.rgb = C_TEXT_MUTED

    add_footer(s9, 9)

    # -------------------------------------------------------------------------
    # SLIDE 10: EVIDENCIAS COMPROMETIDAS (Con Diagrama Voice SOAR)
    # -------------------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    apply_lofi_background(s10)
    add_header(s10, "Evidencias Comprometidas del Proyecto")

    add_card(s10, 0.8, 1.35, 5.6, 2.75, title="EVIDENCIAS DE AVANCE (FASE 1 Y 2)", border_color=C_CYAN)
    tb_e1 = s10.shapes.add_textbox(Inches(0.95), Inches(1.80), Inches(5.3), Inches(2.2))
    tf_e1 = tb_e1.text_frame
    tf_e1.word_wrap = True
    p_e1 = tf_e1.paragraphs[0]
    p_e1.text = (
        "• [F1] Informe de Definición y Topología: Diseño L2/L3, direccionamiento IP y fundamentación.\n"
        "• [F2] Repositorio Git e IaC: Código Python KRONOS, Dockerfiles de Asterisk/DVWA y XML pfSense.\n"
        "• [F2] Manuales Técnicos PDF: Manual pfSense 2.9.0 y Tutorial Paso a Paso con WebGUI mockups."
    )
    p_e1.font.name = "Segoe UI"
    p_e1.font.size = Pt(8.5)
    p_e1.font.color.rgb = C_TEXT_LIGHT

    add_card(s10, 0.8, 4.25, 5.6, 2.65, title="EVIDENCIAS FINALES Y DEMO EN VIVO (FASE 3)", border_color=C_GREEN)
    tb_e2 = s10.shapes.add_textbox(Inches(0.95), Inches(4.70), Inches(5.3), Inches(2.1))
    tf_e2 = tb_e2.text_frame
    tf_e2.word_wrap = True
    p_e2 = tf_e2.paragraphs[0]
    p_e2.text = (
        "• [F3] Demostración en Vivo SOAR: Inyección SQL en DVWA ➔ Bloqueo Netmap ➔ pfctl Kill States ➔ Llamada Gemini Live al CISO.\n"
        "• [F3] Matrices QA y Logs: EVE JSON, tablas snort2c y latencia <1.5s.\n"
        "• [F3] Informe Final Consolidado y Defensa ante Comisión."
    )
    p_e2.font.name = "Segoe UI"
    p_e2.font.size = Pt(8.5)
    p_e2.font.color.rgb = C_TEXT_LIGHT

    # Diagrama Voice SOAR
    voice_path = os.path.join(ASSETS_DIR, "voice_soar_flow.png")
    add_image_pro(s10, voice_path, 6.70, 1.35, 5.833, 3.30)

    # Widget de Telemetría
    add_card(s10, 6.70, 4.85, 5.833, 2.05, title="TELEMETRÍA Y RESPUESTA EN TIEMPO REAL", border_color=C_PURPLE, header_h=0.32)
    tb_tel = s10.shapes.add_textbox(Inches(6.85), Inches(5.25), Inches(5.5), Inches(1.55))
    tf_tel = tb_tel.text_frame
    tf_tel.word_wrap = True
    p_tel = tf_tel.paragraphs[0]
    p_tel.text = (
        "• Protocolo de Voz: Streaming PCM 24kHz dúplex directo con Google Gemini Live.\n"
        "• Interacción Táctica: El CISO escucha el vector mitigado y puede autorizar acciones de respuesta adicionales por voz en lenguaje natural."
    )
    p_tel.font.name = "Segoe UI"
    p_tel.font.size = Pt(8.3)
    p_tel.font.color.rgb = C_TEXT_LIGHT

    add_footer(s10, 10)

    # -------------------------------------------------------------------------
    # SLIDE 11: PLAN DE TRABAJO Y CARTA GANTT (Widescreen 3.66:1)
    # -------------------------------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    apply_lofi_background(s11)
    add_header(s11, "Plan de Trabajo y Carta Gantt (18 Semanas)")

    # Gráfico Gantt Widescreen Proporcional
    gantt_path = os.path.join(ASSETS_DIR, "gantt_fase1_timeline.png")
    add_image_pro(s11, gantt_path, 0.8, 1.35, 11.733, 3.20)

    # 4 Tarjetas de Fases Inferiores
    fases_g = [
        ("FASE 1 (S1-S4)", "Setup pfSense & VLANs 802.1Q (Corp 10, DMZ 20, VoIP 30). Informe EA1.", C_CYAN),
        ("FASE 2A (S5-S10)", "Suricata Netmap IPS, GeoIP pfBlockerNG y HAProxy SSL DMZ.", C_SKY),
        ("FASE 2B (S11-S15)", "Motor KRONOS AST, Asterisk PBX, Gemini Live Voice y Tailscale.", C_AMBER),
        ("FASE 3 (S16-S18)", "Pruebas QA <1.5s, Manuales PDF y Defensa Portafolio APT122.", C_GREEN)
    ]
    card_wg = 2.78
    for idx, (title, desc, col) in enumerate(fases_g):
        x_g = 0.8 + idx * (card_wg + 0.20)
        add_card(s11, x_g, 4.80, card_wg, 2.10, title=title, border_color=col, header_h=0.35)
        tb_g = s11.shapes.add_textbox(Inches(x_g + 0.12), Inches(5.25), Inches(card_wg - 0.24), Inches(1.55))
        tf_g = tb_g.text_frame
        tf_g.word_wrap = True
        p = tf_g.paragraphs[0]
        p.text = desc
        p.font.name = "Segoe UI"
        p.font.size = Pt(8.5)
        p.font.color.rgb = C_TEXT_LIGHT

    add_footer(s11, 11)

    # -------------------------------------------------------------------------
    # SLIDE 12: CIERRE, PRÓXIMOS PASOS (FASE 2) Y PREGUNTAS
    # -------------------------------------------------------------------------
    s12 = prs.slides.add_slide(blank_layout)
    apply_lofi_background(s12)
    add_header(s12, "Cierre, Próximos Pasos (Fase 2) y Preguntas")

    # Tarjeta Izquierda: Próximos Pasos
    add_card(s12, 0.8, 1.35, 5.7, 5.55, title="PRÓXIMOS PASOS HACIA FASE 2", border_color=C_CYAN)
    tb_nxt = s12.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(5.3), Inches(4.9))
    tf_nxt = tb_nxt.text_frame
    tf_nxt.word_wrap = True
    nxt_text = (
        "1. Despliegue de Suricata Inline Netmap:\n"
        "• Configurar dropsid.conf para conversión automática de alertas a drops en hardware ring-buffer.\n\n"
        "2. Hardening de HAProxy en DMZ:\n"
        "• Parametrizar Stick-Tables para rate-limiting dinámico anti-fuzzing L7.\n\n"
        "3. Desarrollo del Motor KRONOS (Python 3.12):\n"
        "• Implementar AST parser heurístico y llamadas subprocess a pfctl -k.\n\n"
        "4. Pruebas de Integración Asterisk & Gemini Live:\n"
        "• Validar streaming de audio bidireccional PCM 24kHz y timbrado a PJSIP/1001."
    )
    p = tf_nxt.paragraphs[0]
    p.text = nxt_text
    p.font.name = "Segoe UI"
    p.font.size = Pt(9.2)
    p.font.color.rgb = C_TEXT_LIGHT

    # Tarjeta Derecha: Espacio de Preguntas
    add_card(s12, 6.833, 1.35, 5.7, 5.55, title="ESPACIO DE PREGUNTAS Y CONSULTAS", border_color=C_GREEN)
    
    # Emblema Centrado (Dark Purple Seamless)
    add_image_pro(s12, logo_path, 8.883, 1.95, 1.6, 1.6)

    tb_q = s12.shapes.add_textbox(Inches(7.033), Inches(3.70), Inches(5.3), Inches(3.0))
    tf_q = tb_q.text_frame
    tf_q.word_wrap = True
    
    p_q1 = tf_q.paragraphs[0]
    p_q1.text = "KRONOS SENTINEL\n"
    p_q1.font.name = "Segoe UI"
    p_q1.font.size = Pt(14)
    p_q1.font.bold = True
    p_q1.font.color.rgb = C_TEXT_WHITE

    p_q2 = tf_q.add_paragraph()
    p_q2.text = "“Innovación en Ciberdefensa: Desacoplando la contención atómica en kernel de la comunicación estratégica de voz en tiempo real.”\n\n¡Muchas gracias por su atención!\nQuedamos a disposición de la Comisión Evaluadora de Duoc UC."
    p_q2.font.name = "Segoe UI"
    p_q2.font.size = Pt(9.0)
    p_q2.font.color.rgb = C_TEXT_LIGHT

    add_footer(s12, 12)

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    target_files = [
        "Presentacion_Proyecto_APT_Fase1_KRONOS_SENTINEL_FINAL.pptx",
        "Presentacion_Proyecto_APT_Fase1_KRONOS_SENTINEL_LOFI.pptx",
        "Presentacion_Proyecto_APT_Fase1_KRONOS_SENTINEL.pptx"
    ]
    saved_paths = []
    for fname in target_files:
        fpath = os.path.join(OUTPUT_DIR, fname)
        try:
            prs.save(fpath)
            saved_paths.append(fpath)
            print(f"[OK] PowerPoint saved successfully: {fpath}")
        except PermissionError:
            print(f"[NOTE] '{fname}' is currently open in PowerPoint (skipped write).")

if __name__ == "__main__":
    build_presentation()
